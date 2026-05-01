"""
rag/document_loader.py
======================
Modul untuk memuat dokumen dari berbagai format:
- PDF (.pdf)
- PowerPoint Slides (.pptx)
- Word Document (.docx)
- Text file (.txt)
"""

import os
from pathlib import Path
from typing import List

from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
)
from langchain_community.document_loaders import UnstructuredPowerPointLoader


class DocumentLoader:
    """
    Kelas untuk memuat dan memproses dokumen dari berbagai format file.
    Mendukung: PDF, PPTX, DOCX, TXT
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf": "PDF",
        ".pptx": "PowerPoint",
        ".ppt": "PowerPoint (lama)",
        ".docx": "Word Document",
        ".txt": "Text File",
    }

    def __init__(self):
        self.loaded_files = []

    def load_file(self, file_path: str) -> List[Document]:
        """
        Memuat satu file dan mengembalikan list Document.

        Args:
            file_path: Path ke file yang akan dimuat

        Returns:
            List of LangChain Document objects
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        if not path.exists():
            raise FileNotFoundError(f"File tidak ditemukan: {file_path}")

        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Format file '{ext}' tidak didukung. "
                f"Format yang didukung: {list(self.SUPPORTED_EXTENSIONS.keys())}"
            )

        print(f"[📄] Memuat {self.SUPPORTED_EXTENSIONS[ext]}: {path.name}")

        try:
            if ext == ".pdf":
                docs = self._load_pdf(file_path)
            elif ext in (".pptx", ".ppt"):
                docs = self._load_pptx(file_path)
            elif ext == ".docx":
                docs = self._load_docx(file_path)
            elif ext == ".txt":
                docs = self._load_txt(file_path)
            else:
                docs = []

            # Tambahkan metadata sumber file
            for doc in docs:
                doc.metadata["source_file"] = path.name
                doc.metadata["file_type"] = self.SUPPORTED_EXTENSIONS[ext]

            self.loaded_files.append(path.name)
            print(f"   ✅ Berhasil dimuat: {len(docs)} halaman/bagian")
            return docs

        except Exception as e:
            print(f"   ❌ Gagal memuat {path.name}: {e}")
            raise

    def load_directory(self, directory_path: str) -> List[Document]:
        """
        Memuat semua dokumen yang didukung dari sebuah folder.

        Args:
            directory_path: Path ke folder dokumen

        Returns:
            List of semua Document yang berhasil dimuat
        """
        directory = Path(directory_path)

        if not directory.exists():
            raise FileNotFoundError(f"Folder tidak ditemukan: {directory_path}")

        all_docs = []
        files_found = list(directory.rglob("*"))
        supported_files = [
            f for f in files_found
            if f.is_file() and f.suffix.lower() in self.SUPPORTED_EXTENSIONS
        ]

        if not supported_files:
            print(f"⚠️  Tidak ada file yang didukung di folder: {directory_path}")
            print(f"   Format yang didukung: {list(self.SUPPORTED_EXTENSIONS.keys())}")
            return []

        print(f"\n[📁] Ditemukan {len(supported_files)} file di '{directory_path}':")
        for f in supported_files:
            print(f"   - {f.name}")

        print()
        for file_path in supported_files:
            try:
                docs = self.load_file(str(file_path))
                all_docs.extend(docs)
            except Exception as e:
                print(f"   ⚠️  Melewati {file_path.name}: {e}")

        print(f"\n[✅] Total dokumen dimuat: {len(all_docs)} bagian dari {len(supported_files)} file\n")
        return all_docs

    def _load_pdf(self, file_path: str) -> List[Document]:
        """Memuat file PDF menggunakan PyPDFLoader."""
        loader = PyPDFLoader(file_path)
        return loader.load()

    def _load_pptx(self, file_path: str) -> List[Document]:
        """Memuat file PowerPoint menggunakan UnstructuredPowerPointLoader."""
        try:
            loader = UnstructuredPowerPointLoader(file_path)
            return loader.load()
        except Exception:
            # Fallback: baca langsung dengan python-pptx
            return self._load_pptx_manual(file_path)

    def _load_pptx_manual(self, file_path: str) -> List[Document]:
        """Fallback loader untuk PPTX menggunakan python-pptx langsung."""
        from pptx import Presentation

        prs = Presentation(file_path)
        docs = []

        for i, slide in enumerate(prs.slides):
            text_parts = []

            # Ambil teks dari semua shape di slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

            if text_parts:
                content = "\n".join(text_parts)
                doc = Document(
                    page_content=content,
                    metadata={"page": i + 1, "slide_number": i + 1}
                )
                docs.append(doc)

        return docs if docs else [Document(
            page_content="[Slide kosong atau tidak ada teks]",
            metadata={"page": 1}
        )]

    def _load_docx(self, file_path: str) -> List[Document]:
        """Memuat file Word Document."""
        loader = Docx2txtLoader(file_path)
        return loader.load()

    def _load_txt(self, file_path: str) -> List[Document]:
        """Memuat file teks biasa."""
        loader = TextLoader(file_path, encoding="utf-8")
        return loader.load()

    def get_loaded_files(self) -> List[str]:
        """Mengembalikan list nama file yang sudah dimuat."""
        return self.loaded_files.copy()
